"""Web UI server for the Software Factory demo.

A dependency-free (stdlib only) HTTP server that:
  * serves the single-page UI (index.html)
  * streams the factory workflow as Server-Sent Events (/api/run)
  * builds a .pptx from the report text on demand (/api/report-pptx)
  * reports CloudWatch logging status (/api/status)

Every workflow step is mirrored to CloudWatch Logs by the director, so all
invocations are observable in the AWS console while the UI visualizes them live.

Run:
    python -m webui.server
    # then open http://localhost:8080
"""

from __future__ import annotations

import json
import os
import urllib.parse
from http.server import BaseHTTPRequestHandler, ThreadingHTTPServer
from pathlib import Path

from webui.core import CONFIG, CLOUDWATCH, LOG_GROUP, TENANTS
from webui import evals, live_invoker, pipeline

WEBUI_DIR = Path(__file__).parent
HOST = os.getenv("WEBUI_HOST", "127.0.0.1")
PORT = int(os.getenv("WEBUI_PORT", "8090"))


def _build_pptx(title: str, tenant: str, text: str) -> bytes:
    """Render report text into a simple .pptx and return the bytes."""
    import io
    from pptx import Presentation
    from pptx.util import Pt

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = f"Tenant: {tenant}" if tenant else ""

    body_slide = prs.slides.add_slide(prs.slide_layouts[1])
    body_slide.shapes.title.text = "Summary"
    tf = body_slide.placeholders[1].text_frame
    lines = (text or "").strip().splitlines() or ["(no content)"]
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = line[:200]
        para.font.size = Pt(12)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _online_eval_status() -> dict:
    """Best-effort lookup of the deployed online eval config status."""
    try:
        import boto3

        c = boto3.client("bedrock-agentcore-control", region_name=CONFIG.aws_region)
        resp = c.list_online_evaluation_configs()
        key = next((k for k in resp if "onfig" in k and isinstance(resp[k], list)), None)
        configs = resp.get(key, []) if key else []
        if configs:
            return {"available": True, "status": configs[0].get("status", "UNKNOWN"),
                    "count": len(configs)}
        return {"available": False, "status": "none"}
    except Exception as exc:  # noqa: BLE001
        return {"available": False, "status": f"error: {type(exc).__name__}"}

class Handler(BaseHTTPRequestHandler):
    def log_message(self, *args):  # quiet default logging
        pass

    # -- helpers ---------------------------------------------------------
    def _send(self, code: int, body: bytes, content_type: str = "text/plain") -> None:
        self.send_response(code)
        self.send_header("Content-Type", content_type)
        self.send_header("Content-Length", str(len(body)))
        self.end_headers()
        self.wfile.write(body)

    # -- routing ---------------------------------------------------------
    def do_GET(self) -> None:
        parsed = urllib.parse.urlparse(self.path)
        route = parsed.path

        if route in ("/", "/index.html"):
            self._serve_file(WEBUI_DIR / "index.html", "text/html")
        elif route == "/app.js":
            self._serve_file(WEBUI_DIR / "app.js", "application/javascript")
        elif route == "/styles.css":
            self._serve_file(WEBUI_DIR / "styles.css", "text/css")
        elif route == "/api/status":
            self._api_status()
        elif route == "/api/tenants":
            self._api_tenants()
        elif route == "/api/run":
            self._api_run(urllib.parse.parse_qs(parsed.query))
        elif route == "/api/evals":
            self._api_evals(urllib.parse.parse_qs(parsed.query))
        elif route == "/api/pipeline":
            self._api_pipeline()
        else:
            self._send(404, b"not found")

    def do_POST(self) -> None:
        parsed = urllib.parse.urlparse(self.path)
        if parsed.path == "/api/report-pptx":
            self._api_report_pptx()
        else:
            self._send(404, b"not found")

    def _api_report_pptx(self) -> None:
        """Build a .pptx from posted report text (used in LIVE mode where the
        deployed director returns text, not a saved artifact)."""
        length = int(self.headers.get("Content-Length", "0"))
        body = self.rfile.read(length) if length else b"{}"
        try:
            payload = json.loads(body.decode("utf-8"))
        except Exception:
            payload = {}
        title = payload.get("title", "Customer Status Report")
        tenant = payload.get("tenant", "")
        text = payload.get("text", "")
        try:
            data = _build_pptx(title, tenant, text)
        except Exception as exc:  # noqa: BLE001
            self._send(500, f"pptx build failed: {exc}".encode())
            return
        self.send_response(200)
        self.send_header(
            "Content-Type",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
        self.send_header("Content-Disposition", 'attachment; filename="customer_status_report.pptx"')
        self.send_header("Content-Length", str(len(data)))
        self.end_headers()
        self.wfile.write(data)

    # -- endpoints -------------------------------------------------------
    def _serve_file(self, path: Path, content_type: str) -> None:
        if not path.exists():
            self._send(404, b"not found")
            return
        self._send(200, path.read_bytes(), content_type)

    def _api_status(self) -> None:
        payload = {
            "mode": "LIVE (AgentCore)",
            "region": CONFIG.aws_region,
            "model": CONFIG.model_id,
            "cloudwatch_enabled": CLOUDWATCH.enabled,
            "cloudwatch_status": CLOUDWATCH.status,
            "cloudwatch_log_group": LOG_GROUP,
            "live_available": live_invoker.is_available(),
            "online_eval": _online_eval_status(),
        }
        self._send(200, json.dumps(payload).encode(), "application/json")

    def _api_tenants(self) -> None:
        payload = [
            {"id": t.tenant_id, "name": t.display_name,
             "rate_limit": t.rate_limit_per_min}
            for t in TENANTS.values()
        ]
        self._send(200, json.dumps(payload).encode(), "application/json")

    def _api_evals(self, query: dict) -> None:
        which = query.get("evaluators", [",".join(evals.UI_EVALUATORS)])[0]
        evaluators = [e for e in which.split(",") if e]
        days = int(query.get("days", ["1"])[0])
        result = evals.run_ondemand(evaluators, days=days)
        self._send(200, json.dumps(result).encode(), "application/json")

    def _api_pipeline(self) -> None:
        self.send_response(200)
        self.send_header("Content-Type", "text/event-stream")
        self.send_header("Cache-Control", "no-cache")
        self.send_header("Connection", "keep-alive")
        self.end_headers()

        def write_event(obj: dict) -> None:
            self.wfile.write(f"data: {json.dumps(obj)}\n\n".encode())
            self.wfile.flush()

        try:
            for event in pipeline.run_pipeline():
                write_event(event)
            write_event({"stage": "__end__", "status": "done", "detail": ""})
        except BrokenPipeError:
            pass
        except Exception as exc:  # noqa: BLE001
            try:
                write_event({"stage": "score", "status": "error", "detail": str(exc)})
            except Exception:
                pass

    def _api_run(self, query: dict) -> None:
        prompt = query.get("prompt", ["Generate a customer status report"])[0]
        tenant_id = query.get("tenant", ["msp-a"])[0]
        user_id = query.get("user", ["eng-1"])[0]
        session_id = query.get(
            "session", [f"{tenant_id}-{user_id}-uisession-000000000001"]
        )[0]

        # SSE headers
        self.send_response(200)
        self.send_header("Content-Type", "text/event-stream")
        self.send_header("Cache-Control", "no-cache")
        self.send_header("Connection", "keep-alive")
        self.end_headers()

        def write_event(obj: dict) -> None:
            self.wfile.write(f"data: {json.dumps(obj)}\n\n".encode())
            self.wfile.flush()

        try:
            # Drive the DEPLOYED AgentCore director runtime (LIVE only).
            for event in live_invoker.run_live(tenant_id, session_id, user_id, prompt):
                write_event(event.to_dict())
            write_event({"step": "__end__", "status": "done", "title": "", "detail": ""})
        except BrokenPipeError:
            pass
        except Exception as exc:  # noqa: BLE001
            try:
                write_event({"step": "director", "status": "error",
                             "title": "Director agent", "detail": str(exc), "data": {}})
            except Exception:
                pass

def main() -> None:
    candidates = [PORT, 8090, 8188, 8200, 8500, 9000]
    server = None
    bound_port = None
    for port in dict.fromkeys(candidates):
        try:
            server = ThreadingHTTPServer((HOST, port), Handler)
            bound_port = port
            break
        except (PermissionError, OSError):
            continue
    if server is None:
        print("Could not bind any port. Set WEBUI_PORT to a free port.")
        return

    print("\n" + "=" * 60)
    print("  Software Factory - Live UI")
    print("=" * 60)
    print(f"  URL         : http://{HOST}:{bound_port}")
    print(f"  Mode        : LIVE (deployed AgentCore)")
    print(f"  CloudWatch  : {CLOUDWATCH.status}")
    print(f"  Log group   : {LOG_GROUP}")
    print("=" * 60 + "\n")
    try:
        server.serve_forever()
    except KeyboardInterrupt:
        print("\nShutting down.")
        server.shutdown()


if __name__ == "__main__":
    main()
